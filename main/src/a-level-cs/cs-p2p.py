from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
slide_layout = prs.slide_layouts[5]  # Title Only

def add_slide(title, points, icon):
    slide = prs.slides.add_slide(slide_layout)
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True

    # Bullets (left)
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(4))
    tf = left_box.text_frame
    for i, point in enumerate(points):
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(24)
        if i == 0:
            p.font.bold = True

    # Icon (right)
    icon_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.5), Inches(3), Inches(3))
    icon_shape.text = icon
    icon_shape.text_frame.paragraphs[0].font.size = Pt(70)
    icon_shape.text_frame.paragraphs[0].font.bold = True
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = RGBColor(91, 155, 213)

# Slide 1: Title
add_slide("How do computers share files?", [
    "Let’s find out in 5 minutes!",
], "💡")

# Slide 2: The Problem
add_slide("Ever had this happen?", [
    "You can’t log in at school.",
    "AirDrop works, but WiFi doesn’t.",
    "Why is that?",
], "🤔")

# Slide 3: Client-Server
add_slide("Client-Server: School Login", [
    "One big computer (server) in charge.",
    "You ask, it gives access.",
    "Secure, but if it fails—no one gets in!",
], "🖥️")

# Slide 4: Peer-to-Peer
add_slide("Peer-to-Peer: Like AirDrop", [
    "No boss—everyone’s equal.",
    "Share files directly.",
    "Fast, but less secure.",
], "📱")

# Slide 5: Quick Quiz!
add_slide("Which is which?", [
    "1. School login?",
    "2. AirDrop?",
    "3. Minecraft LAN?",
    "Hold up 1 or 2 fingers!",
], "❓")

# Slide 6: Why it matters
add_slide("Why care?", [
    "Pick the right network for the job.",
    "Stay safe online.",
    "You’re the next IT expert!",
], "🔐")

prs.save('Peer_to_Peer_vs_Client_Server_Lesson_Engaging.pptx')