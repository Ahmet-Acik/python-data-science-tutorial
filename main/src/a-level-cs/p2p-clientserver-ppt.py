import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
title_slide_layout = prs.slide_layouts[5]  # Title Only

def add_half_slide(title, content_lines, icon_text=None, image_path=None):
    slide = prs.slides.add_slide(title_slide_layout)
    # Title at the top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True

    # Text on left half
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(4.5)
    height = Inches(4.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    for i, line in enumerate(content_lines):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(26)
        p.font.name = "Arial"
        if i == 0:
            p.font.bold = True
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    # Icon or image on right half
    right_left = Inches(5.5)
    right_top = Inches(1.7)
    right_width = Inches(3)
    right_height = Inches(3)
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, right_left, right_top, right_width, right_height)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_left, right_top, right_width, right_height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(91, 155, 213)
        if icon_text:
            shape.text = icon_text
            shape.text_frame.paragraphs[0].font.size = Pt(60)
            shape.text_frame.paragraphs[0].font.bold = True

# Title slide
slide = prs.slides.add_slide(title_slide_layout)
# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "How do computers talk to each other?"
title_frame.paragraphs[0].font.size = Pt(40)
title_frame.paragraphs[0].font.bold = True
# Subtitle on left
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Exploring Peer-to-Peer and Client-Server Networks"
subtitle_frame.paragraphs[0].font.size = Pt(28)
# Icon on right
icon_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.5), Inches(3), Inches(3))
icon_shape.text = "💻"
icon_shape.text_frame.paragraphs[0].font.size = Pt(70)
icon_shape.text_frame.paragraphs[0].font.bold = True
icon_shape.fill.solid()
icon_shape.fill.fore_color.rgb = RGBColor(91, 155, 213)

# Content slides
add_half_slide("Hook: Streaming vs Sharing", [
    "Imagine one illegal movie download being shared 10 million times — no website involved!",
    "Or your school login not working because a server is down.",
    "Sharing music via AirDrop vs. streaming from Spotify — what's the difference?"
], icon_text="🎧")

add_half_slide("Peer-to-Peer (P2P) – Like Bluetooth Sharing", [
    "No central server or admin",
    "Devices connect directly to share files",
    "Each device acts as both sender and receiver",
    "Examples: Bluetooth, AirDrop, Minecraft LAN, file-sharing apps"
], icon_text="📲")

add_half_slide("Client-Server – Like Streaming or Logging In", [
    "One main server provides services to others (clients)",
    "Devices request files, login access, or websites",
    "Easier to control and manage",
    "Examples: School logins, Google Classroom, Fortnite online"
], icon_text="🌐")

add_half_slide("P2P vs Client-Server – Spot the Differences", [
    "P2P: Shared control, low cost, hard to secure, e.g., Minecraft LAN",
    "Client-Server: Central control, higher cost, secure, e.g., Google Classroom"
], icon_text="⚖️")

add_half_slide("Quick Quiz - Which Model?", [
    "1. Sam shares music via AirDrop → ?",
    "2. Jade logs into school email → ?",
    "3. Tom and Mia play Minecraft in same room → ?",
    "Hold up 1 finger for P2P, 2 for Client-Server!"
], icon_text="❓")

add_half_slide("Why does this matter in real life?", [
    "Helps choose the right network",
    "Schools need security → Client-Server",
    "No internet? Use Peer-to-Peer with friends"
], icon_text="🔐")

prs.save("Peer_to_Peer_vs_Client_Server_Lesson_Engaging.pptx")