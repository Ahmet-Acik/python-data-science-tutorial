import os

pptx_path = '/Users/drahmetacik/Projects/py-np-pd-plt-project/main/src/a-level-cs/client-server-p2p.pptx'
if not os.path.exists(pptx_path):
    print(f"File not found: {pptx_path}")
else:
    from pptx import Presentation

    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i + 1}:")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print("  Text:", shape.text.strip())
        print("-" * 40)

import os
from pptx import Presentation
from pptx.util import Pt

input_pptx = '/Users/drahmetacik/Projects/py-np-pd-plt-project/main/src/a-level-cs/client-server-p2p.pptx'
output_pptx = '/Users/drahmetacik/Projects/py-np-pd-plt-project/main/src/a-level-cs/ppt.analysis.generated.pptx'

MAX_LINES_PER_SLIDE = 10  # Adjust as needed

if not os.path.exists(input_pptx):
    print(f"File not found: {input_pptx}")
else:
    prs_in = Presentation(input_pptx)
    prs_out = Presentation()
    # Remove default first slide
    if len(prs_out.slides) > 0:
        xml_slides = prs_out.slides._sldIdLst
        slides = list(xml_slides)
        for slide in slides:
            xml_slides.remove(slide)

    for slide in prs_in.slides:
        slide_texts = [shape.text.strip() for shape in slide.shapes if
                       hasattr(shape, "text") and shape.text.strip()]
        if not slide_texts:
            slide_texts = [""]
        title = slide_texts[0]
        content_lines = slide_texts[1:] if len(slide_texts) > 1 else []
        # Split content into multiple slides if too long
        for i in range(0, max(1, len(content_lines)), MAX_LINES_PER_SLIDE):
            new_slide = prs_out.slides.add_slide(prs_out.slide_layouts[1])
            new_slide.shapes.title.text = title if i == 0 else f"{title} (cont.)"
            content_shape = new_slide.placeholders[1]
            lines = content_lines[i:i + MAX_LINES_PER_SLIDE]
            content_shape.text = "\n".join(lines)
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
            content_shape.text_frame.word_wrap = True
            # If no content, clear placeholder
            if not lines:
                content_shape.text = ""
    prs_out.save(output_pptx)
    print(f"Presentation created at {output_pptx}")
